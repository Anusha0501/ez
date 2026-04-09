"""
PowerPoint generation utilities.
"""

import os
import logging
from typing import List, Dict, Any, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import io
import base64

from core.models import Slide, SlideLayout, VisualElement, ChartData, ChartType


class PPTXUtils:
    """Utilities for creating PowerPoint presentations."""
    
    def __init__(self):
        self.logger = logging.getLogger("pptx_utils")
        plt.style.use('default')
    
    def create_presentation(self, template_path: Optional[str] = None) -> Presentation:
        """Create a new presentation, optionally from a template."""
        if template_path and os.path.exists(template_path):
            return Presentation(template_path)
        else:
            return Presentation()
    
    def add_slide(self, prs: Presentation, layout_index: int = 6) -> Any:
        """Add a slide with the specified layout."""
        slide_layout = prs.slide_layouts[layout_index]
        return prs.slides.add_slide(slide_layout)
    
    def add_title_slide(self, prs: Presentation, title: str, subtitle: str = "") -> Any:
        """Add a title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        if subtitle and slide.placeholders[1:]:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
        
        return slide
    
    def add_content_slide(self, prs: Presentation, title: str, content: List[str]) -> Any:
        """Add a content slide with bullet points."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        if slide.placeholders[1:]:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for i, point in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
        
        return slide
    
    def add_two_content_slide(self, prs: Presentation, title: str, 
                             left_content: List[str], right_content: List[str]) -> Any:
        """Add a slide with two content columns."""
        slide_layout = prs.slide_layouts[3]  # Two content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Left content
        if slide.placeholders[1:]:
            left_shape = slide.placeholders[1]
            left_frame = left_shape.text_frame
            left_frame.clear()
            
            for i, point in enumerate(left_content):
                if i == 0:
                    p = left_frame.paragraphs[0]
                else:
                    p = left_frame.add_paragraph()
                p.text = point
                p.level = 0
        
        # Right content
        if slide.placeholders[2:]:
            right_shape = slide.placeholders[2]
            right_frame = right_shape.text_frame
            right_frame.clear()
            
            for i, point in enumerate(right_content):
                if i == 0:
                    p = right_frame.paragraphs[0]
                else:
                    p = right_frame.add_paragraph()
                p.text = point
                p.level = 0
        
        return slide
    
    def add_chart_slide(self, prs: Presentation, title: str, chart_data: ChartData) -> Any:
        """Add a slide with a chart."""
        slide_layout = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Create chart image
        chart_image = self._create_chart_image(chart_data)
        if chart_image:
            # Add the chart image to the slide
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)
            
            slide.shapes.add_picture(chart_image, left, top, width, height)
        
        return slide
    
    def _create_chart_image(self, chart_data: ChartData) -> Optional[Any]:
        """Create a chart image from chart data."""
        try:
            fig, ax = plt.subplots(figsize=(10, 6))
            
            if chart_data.chart_type == ChartType.BAR:
                self._create_bar_chart(ax, chart_data)
            elif chart_data.chart_type == ChartType.PIE:
                self._create_pie_chart(ax, chart_data)
            elif chart_data.chart_type == ChartType.LINE:
                self._create_line_chart(ax, chart_data)
            else:
                self.logger.warning(f"Unsupported chart type: {chart_data.chart_type}")
                return None
            
            ax.set_title(chart_data.title, fontsize=16, fontweight='bold')
            plt.tight_layout()
            
            # Save to bytes
            img_buffer = io.BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            return img_buffer
            
        except Exception as e:
            self.logger.error(f"Error creating chart: {str(e)}")
            return None
    
    def _create_bar_chart(self, ax, chart_data: ChartData):
        """Create a bar chart."""
        if len(chart_data.datasets) == 0:
            return
        
        dataset = chart_data.datasets[0]
        x = range(len(chart_data.labels))
        
        bars = ax.bar(x, dataset['data'], color=chart_data.colors[0] if chart_data.colors else '#1f77b4')
        
        ax.set_xticks(x)
        ax.set_xticklabels(chart_data.labels, rotation=45, ha='right')
        ax.set_ylabel(dataset.get('label', 'Value'))
        
        # Add value labels on bars
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height,
                   f'{height:.1f}', ha='center', va='bottom')
    
    def _create_pie_chart(self, ax, chart_data: ChartData):
        """Create a pie chart."""
        if len(chart_data.datasets) == 0:
            return
        
        dataset = chart_data.datasets[0]
        colors = chart_data.colors if chart_data.colors else plt.cm.Set3(range(len(chart_data.labels)))
        
        wedges, texts, autotexts = ax.pie(
            dataset['data'], 
            labels=chart_data.labels,
            colors=colors,
            autopct='%1.1f%%',
            startangle=90
        )
        
        # Ensure pie chart is circular
        ax.axis('equal')
    
    def _create_line_chart(self, ax, chart_data: ChartData):
        """Create a line chart."""
        x = range(len(chart_data.labels))
        
        for i, dataset in enumerate(chart_data.datasets):
            color = chart_data.colors[i] if chart_data.colors and i < len(chart_data.colors) else None
            ax.plot(x, dataset['data'], marker='o', label=dataset.get('label', f'Series {i+1}'), color=color)
        
        ax.set_xticks(x)
        ax.set_xticklabels(chart_data.labels, rotation=45, ha='right')
        ax.set_ylabel('Value')
        ax.legend()
        ax.grid(True, alpha=0.3)
    
    def add_visual_element(self, slide, visual_element: VisualElement):
        """Add a visual element to a slide."""
        try:
            if visual_element.type == "flow":
                self._add_flow_diagram(slide, visual_element)
            elif visual_element.type == "timeline":
                self._add_timeline(slide, visual_element)
            elif visual_element.type == "grid":
                self._add_grid_layout(slide, visual_element)
            else:
                self.logger.warning(f"Unsupported visual element type: {visual_element.type}")
                
        except Exception as e:
            self.logger.error(f"Error adding visual element: {str(e)}")
    
    def _add_flow_diagram(self, slide, visual_element: VisualElement):
        """Add a flow diagram to a slide."""
        # Create a simple flow diagram using shapes
        content = visual_element.content
        layout = visual_element.layout
        
        left = Inches(layout.get('left', 1))
        top = Inches(layout.get('top', 2))
        width = Inches(layout.get('width', 8))
        height = Inches(layout.get('height', 5))
        
        # Add rectangles for each step
        box_width = width / len(content)
        
        for i, item in enumerate(content):
            x = left + i * box_width
            y = top
            
            # Add rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, box_width - Inches(0.2), height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(70, 130, 180)  # Steel blue
            shape.line.color.rgb = RGBColor(0, 0, 0)
            
            # Add text
            text_frame = shape.text_frame
            text_frame.text = item.get('text', f"Step {i+1}")
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            # Add arrow if not last item
            if i < len(content) - 1:
                arrow_x = x + box_width - Inches(0.2)
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, arrow_x, y + height/2 - Inches(0.2), 
                    Inches(0.4), Inches(0.4)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    def _add_timeline(self, slide, visual_element: VisualElement):
        """Add a timeline to a slide."""
        content = visual_element.content
        layout = visual_element.layout
        
        left = Inches(layout.get('left', 1))
        top = Inches(layout.get('top', 3))
        width = Inches(layout.get('width', 8))
        height = Inches(layout.get('height', 4))
        
        # Draw timeline line
        line_y = top + height / 2
        
        # Add timeline points
        point_spacing = width / (len(content) + 1)
        
        for i, item in enumerate(content):
            x = left + (i + 1) * point_spacing
            
            # Add circle for point
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, x - Inches(0.1), line_y - Inches(0.1), 
                Inches(0.2), Inches(0.2)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(220, 20, 60)  # Crimson
            circle.line.color.rgb = RGBColor(0, 0, 0)
            
            # Add text label
            text_box = slide.shapes.add_textbox(
                x - Inches(0.5), line_y + Inches(0.2), Inches(1), Inches(0.5)
            )
            text_frame = text_box.text_frame
            text_frame.text = item.get('text', f"Point {i+1}")
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_grid_layout(self, slide, visual_element: VisualElement):
        """Add a grid layout to a slide."""
        content = visual_element.content
        layout = visual_element.layout
        
        cols = layout.get('columns', 2)
        rows = layout.get('rows', 2)
        
        left = Inches(layout.get('left', 1))
        top = Inches(layout.get('top', 2))
        width = Inches(layout.get('width', 8))
        height = Inches(layout.get('height', 5))
        
        cell_width = width / cols
        cell_height = height / rows
        
        for i, item in enumerate(content):
            row = i // cols
            col = i % cols
            
            if row >= rows:
                break
            
            x = left + col * cell_width
            y = top + row * cell_height
            
            # Add rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x + Inches(0.05), y + Inches(0.05), 
                cell_width - Inches(0.1), cell_height - Inches(0.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
            shape.line.color.rgb = RGBColor(100, 100, 100)
            
            # Add text
            text_frame = shape.text_frame
            text_frame.text = item.get('text', f"Item {i+1}")
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
    
    def save_presentation(self, prs: Presentation, output_path: str):
        """Save the presentation to a file."""
        try:
            prs.save(output_path)
            self.logger.info(f"Presentation saved to: {output_path}")
        except Exception as e:
            self.logger.error(f"Error saving presentation: {str(e)}")
            raise
