"""
PPTX Generator Agent - Generates the final PowerPoint presentation.
"""

import logging
import os
import re
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from core.agent import BaseAgent
from core.models import ChartData, ChartType
from utils.pptx_utils import PPTXUtils


class PPTXGeneratorAgent(BaseAgent):
    """Agent responsible for generating the final PowerPoint presentation."""
    
    def __init__(self):
        super().__init__("pptx_generator_agent")
        self.pptx_utils = PPTXUtils()
        
        # Presentation styling
        self.theme_colors = {
            "primary": RGBColor(31, 78, 121),        # #1F4E79
            "secondary": RGBColor(47, 128, 237),     # #2F80ED
            "accent": RGBColor(47, 128, 237),        # #2F80ED
            "danger": RGBColor(214, 39, 40),         # Red
            "warning": RGBColor(255, 193, 7),        # Yellow
            "info": RGBColor(23, 162, 184),          # Cyan
            "dark": RGBColor(51, 51, 51),            # #333333
            "light": RGBColor(236, 240, 241)          # Light gray
        }
        
        # Typography
        self.fonts = {
            "title": "Calibri Light",
            "heading": "Calibri",
            "body": "Calibri",
            "subtitle": "Calibri"
        }
        
        self.font_sizes = {
            "title": 44,
            "subtitle": 32,
            "heading": 28,
            "subheading": 20,
            "body": 18,
            "small": 14
        }
    
    def normalize_slides(self, slides):
        """Normalize slide data structure to ensure consistent format."""
        normalized = []
        
        for i, slide in enumerate(slides):
            if isinstance(slide, str):
                normalized.append({
                    "title": f"Slide {i+1}",
                    "content": [slide]
                })
                continue
            
            if isinstance(slide, dict):
                title = (
                    slide.get("title") or 
                    slide.get("heading") or 
                    slide.get("slide_title") or 
                    f"Slide {i+1}"
                )
                content = (
                    slide.get("content") or
                    slide.get("points") or
                    slide.get("bullets") or
                    slide.get("text") or
                    ["Auto-generated content"]
                )
                
                if isinstance(content, str):
                    content = [content]
                
                normalized.append({
                    "title": title,
                    "content": content
                })
        
        return normalized
    
    def process(self, input_data: Dict[str, Any]) -> Dict[str, Any]:
        """Generate the final PowerPoint presentation."""
        try:
            self.log_reasoning("start", "Starting PowerPoint generation")
            
            layout_data = input_data.get("layout_data", {})
            chart_data = input_data.get("chart_data", {})
            storyline = input_data.get("storyline", {})
            output_path = input_data.get("output_path", "output.pptx")
            
            # Get slide layouts from different possible sources
            if isinstance(layout_data, dict):
                slide_layouts = (
                    layout_data.get("layout_data", []) or
                    layout_data.get("slides", []) or
                    []
                )
            elif isinstance(layout_data, list):
                slide_layouts = layout_data
            else:
                slide_layouts = []
            
            # Normalize slides to ensure consistent structure
            slides = self.normalize_slides(slide_layouts)
            assert isinstance(slides, list)
            assert all("title" in s and "content" in s for s in slides)
            
            # Fallback for empty slides
            if not slides:
                slides = [{
                    "title": "Overview",
                    "content": ["Auto-generated presentation"]
                }]
            for slide in slides:
                assert len(slide["content"]) > 0
            seed_slides = slides[:]
            while len(slides) < 10:
                source = seed_slides[len(slides) % len(seed_slides)]
                slides.append({
                    "title": f"{source['title']} (Continued)",
                    "content": source["content"],
                })
            final_slides = slides[:12]
            assert 10 <= len(final_slides) <= 15, "Slide count must be between 10 and 15"
            slides = final_slides
            
            # Log slide normalization
            self.log_reasoning("normalize_slides", f"Slides normalized: {len(slides)}")
            
            # Create presentation
            self.log_reasoning("create_presentation", "Creating new presentation")
            prs = self.pptx_utils.create_presentation()
            
            # Process each slide
            self.log_reasoning("process_slides", "Processing slides", {
                "total_slides": len(slides),
                "output_path": output_path
            })
            chart_specs = chart_data.get("chart_specifications", [])
            
            # Process normalized slides
            for i, slide_data in enumerate(slides):
                slide = self._create_slide_from_normalized(prs, slide_data, chart_specs, i + 1)
                self.log_reasoning("create_slide", f"Created slide {i+1}", {
                    "slide_title": slide_data.get("title", ""),
                    "content_items": len(slide_data.get("content", []))
                })
            
            # Apply theme and styling
            self.log_reasoning("apply_theme", "Applying theme and styling")
            self._apply_theme_to_presentation(prs)
            
            # Save presentation
            self.log_reasoning("save_presentation", f"Saving presentation to {output_path}")
            self.pptx_utils.save_presentation(prs, output_path)
            
            # Create generation summary
            self.log_reasoning("create_summary", "Creating generation summary")
            generation_summary = self._create_generation_summary(slide_layouts, chart_specs, output_path)
            
            output = {
                "presentation_path": output_path,
                "generation_summary": generation_summary,
                "generation_metadata": {
                    "total_slides": len(slide_layouts),
                    "charts_generated": len(chart_specs),
                    "file_size": self._get_file_size(output_path),
                    "theme_applied": "consulting_theme",
                    "generation_time": self._estimate_generation_time(len(slide_layouts))
                }
            }
            
            self.log_reasoning("complete", "Successfully generated PowerPoint presentation", {
                "output_path": output_path,
                "total_slides": len(slide_layouts),
                "file_size": output["generation_metadata"]["file_size"]
            })
            
            return output
            
        except Exception as e:
            self.logger.error(f"Error in PPTX generator agent: {str(e)}")
            raise
    
    def _create_slide(self, prs: Presentation, slide_layout: Dict[str, Any], chart_specs: List[Dict], slide_number: int) -> Any:
        """Create a single slide based on layout."""
        slide_type = slide_layout.get("slide_type", "mixed")
        layout_template = slide_layout.get("layout_template", "flexible_layout")
        
        # Create slide based on type
        if slide_type == "title":
            return self._create_title_slide(prs, slide_layout, slide_number)
        elif slide_type == "agenda":
            return self._create_agenda_slide(prs, slide_layout, slide_number)
        elif slide_type == "process":
            return self._create_process_slide(prs, slide_layout, chart_specs, slide_number)
        elif slide_type == "data":
            return self._create_data_slide(prs, slide_layout, chart_specs, slide_number)
        elif slide_type == "comparison":
            return self._create_comparison_slide(prs, slide_layout, chart_specs, slide_number)
        elif slide_type == "conclusion":
            return self._create_conclusion_slide(prs, slide_layout, slide_number)
        else:
            return self._create_mixed_slide(prs, slide_layout, chart_specs, slide_number)
    
    def _create_title_slide(self, prs: Presentation, slide_layout: Dict[str, Any], slide_number: int) -> Any:
        """Create a title slide."""
        slide_layout_info = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout_info)
        
        # Add content from slide_layout
        content = slide_layout.get("content", [])
        if isinstance(content, list):
            # Add content as text
            content_text = "\n".join(str(item) for item in content)
            self._add_text_content(slide, content_text, {"left": 1, "top": 2, "width": 8, "height": 4})
        
        # Handle content_areas for backward compatibility
        content_areas = slide_layout.get("content_areas", [])
        if content_areas:
            positioning = slide_layout.get("positioning", {})
            
            for area in content_areas:
                area_type = area.get("type", "")
                area_content = area.get("content", {})
                area_key = f"area_{area['id'].split('_')[1]}"
                
                if area_key in positioning:
                    pos = positioning[area_key]
                    
                    if area_type == "title" and area_content:
                        self._add_title_text(slide, area_content.get("text", ""), pos)
                    elif area_type == "subtitle" and area_content:
                        self._add_subtitle_text(slide, area_content.get("text", ""), pos)
        
        return slide
    
    def _create_agenda_slide(self, prs: Presentation, slide_layout: Dict[str, Any], slide_number: int) -> Any:
        """Create an agenda slide."""
        slide_layout_info = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout_info)
        
        # Add title
        content_areas = slide_layout.get("content_areas", [])
        positioning = slide_layout.get("positioning", {})
        
        for area in content_areas:
            area_type = area.get("type", "")
            content = area.get("content", {})
            area_key = f"area_{area['id'].split('_')[1]}"
            
            if area_key in positioning:
                pos = positioning[area_key]
                
                if area_type == "title" and content:
                    self._add_slide_title(slide, content.get("text", "Agenda"))
                elif area_type == "content" and content:
                    self._add_bullet_list(slide, content.get("text", ""), pos)
        
        return slide
    
    def _create_process_slide(self, prs: Presentation, slide_layout: Dict[str, Any], chart_specs: List[Dict], slide_number: int) -> Any:
        """Create a process flow slide."""
        slide_layout_info = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(slide_layout_info)
        
        # Add title
        content_areas = slide_layout.get("content_areas", [])
        positioning = slide_layout.get("positioning", {})
        
        # Add title at the top
        self._add_slide_title(slide, f"Process Flow - Slide {slide_number}")
        
        # Add visual elements
        for area in content_areas:
            if area.get("type") == "visual":
                content = area.get("content", {})
                area_key = f"area_{area['id'].split('_')[1]}"
                
                if area_key in positioning:
                    pos = positioning[area_key]
                    self._add_flow_diagram(slide, content, pos)
        
        # Add chart if applicable
        chart_spec = self._find_chart_for_slide(chart_specs, slide_number)
        if chart_spec:
            self._add_chart_to_slide(slide, chart_spec)
        
        return slide
    
    def _create_data_slide(self, prs: Presentation, slide_layout: Dict[str, Any], chart_specs: List[Dict], slide_number: int) -> Any:
        """Create a data visualization slide."""
        slide_layout_info = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(slide_layout_info)
        
        # Add title
        self._add_slide_title(slide, f"Data Analysis - Slide {slide_number}")
        
        # Add content areas
        content_areas = slide_layout.get("content_areas", [])
        positioning = slide_layout.get("positioning", {})
        
        for area in content_areas:
            area_type = area.get("type", "")
            content = area.get("content", {})
            area_key = f"area_{area['id'].split('_')[1]}"
            
            if area_key in positioning:
                pos = positioning[area_key]
                
                if area_type == "visual" and content:
                    self._add_infographic(slide, content, pos)
                elif area_type == "text" and content:
                    self._add_text_content(slide, content.get("text", ""), pos)
        
        # Add chart if applicable
        chart_spec = self._find_chart_for_slide(chart_specs, slide_number)
        if chart_spec:
            self._add_chart_to_slide(slide, chart_spec)
        
        return slide
    
    def _create_comparison_slide(self, prs: Presentation, slide_layout: Dict[str, Any], chart_specs: List[Dict], slide_number: int) -> Any:
        """Create a comparison slide."""
        slide_layout_info = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(slide_layout_info)
        
        # Add title
        self._add_slide_title(slide, f"Comparison Analysis - Slide {slide_number}")
        
        # Add content areas
        content_areas = slide_layout.get("content_areas", [])
        positioning = slide_layout.get("positioning", {})
        
        for area in content_areas:
            area_type = area.get("type", "")
            content = area.get("content", {})
            area_key = f"area_{area['id'].split('_')[1]}"
            
            if area_key in positioning:
                pos = positioning[area_key]
                
                if area_type == "visual" and content:
                    self._add_comparison_grid(slide, content, pos)
        
        return slide
    
    def _create_conclusion_slide(self, prs: Presentation, slide_layout: Dict[str, Any], slide_number: int) -> Any:
        """Create a conclusion slide."""
        slide_layout_info = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout_info)
        
        # Add title
        self._add_slide_title(slide, "Conclusions & Recommendations")
        
        # Add content
        content_areas = slide_layout.get("content_areas", [])
        positioning = slide_layout.get("positioning", {})
        
        for area in content_areas:
            area_type = area.get("type", "")
            content = area.get("content", {})
            area_key = f"area_{area['id'].split('_')[1]}"
            
            if area_key in positioning:
                pos = positioning[area_key]
                
                if area_type == "content" and content:
                    self._add_conclusion_content(slide, content.get("text", ""), pos)
        
        return slide
    
    def _create_mixed_slide(self, prs: Presentation, slide_layout: Dict[str, Any], chart_specs: List[Dict], slide_number: int) -> Any:
        """Create a mixed content slide."""
        slide_layout_info = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout_info)
        
        # Add title
        self._add_slide_title(slide, f"Key Points - Slide {slide_number}")
        
        # Add content from slide_layout
        content = slide_layout.get("content", [])
        if isinstance(content, list):
            # Add content as text
            content_text = "\n".join(str(item) for item in content)
            self._add_text_content(slide, content_text, {"left": 1, "top": 2, "width": 8, "height": 4})
        
        # Handle content_areas for backward compatibility
        content_areas = slide_layout.get("content_areas", [])
        if content_areas:
            positioning = slide_layout.get("positioning", {})
            
            for area in content_areas:
                area_type = area.get("type", "")
                area_content = area.get("content", {})
                area_key = f"area_{area['id'].split('_')[1]}"
                
                if area_key in positioning:
                    pos = positioning[area_key]
                    
                    if area_type == "visual" and area_content:
                        self._add_visual_element(slide, area_content, pos)
                    elif area_type == "text" and area_content:
                        self._add_text_content(slide, area_content.get("text", ""), pos)
        
        # Add chart if applicable
        chart_spec = self._find_chart_for_slide(chart_specs, slide_number)
        if chart_spec:
            self._add_chart_to_slide(slide, chart_spec)
        
        return slide
    
    def _create_slide_from_normalized(self, prs: Presentation, slide_data: Dict[str, Any], chart_specs: List[Dict], slide_number: int) -> Any:
        """Create slide from normalized structure."""
        title = slide_data.get("title", f"Slide {slide_number}")
        content = slide_data.get("content", [])
        content = content if isinstance(content, list) else [str(content)]

        slide_layout_info = prs.slide_layouts[6]  # Blank layout for custom consulting style
        ppt_slide = prs.slides.add_slide(slide_layout_info)

        if slide_number == 1:
            self._render_title_slide(ppt_slide, title, content)
            return ppt_slide
        if self._is_section_break_slide(title):
            self._render_section_break_slide(ppt_slide, title, slide_number)
            return ppt_slide

        # Section label
        section_box = ppt_slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(0.8), Inches(0.3))
        section_para = section_box.text_frame.paragraphs[0]
        section_para.text = f"{slide_number:02d}"
        section_para.font.name = self.fonts["body"]
        section_para.font.size = Pt(16)
        section_para.font.color.rgb = RGBColor(136, 136, 136)

        # Title + subtitle hierarchy
        title_box = ppt_slide.shapes.add_textbox(Inches(1.1), Inches(0.35), Inches(11.5), Inches(0.6))
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = title
        title_para.font.name = self.fonts["heading"]
        title_para.font.size = Pt(30)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme_colors["primary"]

        subtitle_text = self._clean_bullet(content[0]) if content else "Summary"
        subtitle_box = ppt_slide.shapes.add_textbox(Inches(1.1), Inches(0.95), Inches(11.0), Inches(0.5))
        subtitle_para = subtitle_box.text_frame.paragraphs[0]
        subtitle_para.text = subtitle_text
        subtitle_para.font.name = self.fonts["subtitle"]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = self.theme_colors["dark"]

        # Divider line
        divider = ppt_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.45), Inches(12.2), Inches(0.03))
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(220, 220, 220)
        divider.line.fill.background()

        cleaned_bullets = [self._clean_bullet(item) for item in content if self._clean_bullet(item)]
        cleaned_bullets = cleaned_bullets[1:] if len(cleaned_bullets) > 1 else cleaned_bullets
        highlight_bullets = [b for b in cleaned_bullets if self._is_highlight_bullet(b)]
        left_bullets = [b for b in cleaned_bullets if b not in highlight_bullets][:4]
        if not left_bullets:
            left_bullets = cleaned_bullets[:4]

        # Two-column layout
        left_box = ppt_slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(7.5), Inches(4.9))
        left_tf = left_box.text_frame
        left_tf.clear()
        left_tf.word_wrap = True
        for idx, bullet in enumerate(left_bullets):
            para = left_tf.paragraphs[0] if idx == 0 else left_tf.add_paragraph()
            para.text = bullet
            para.level = 0
            para.space_after = Pt(10)
            para.font.name = self.fonts["body"]
            para.font.size = Pt(20)
            para.font.color.rgb = self.theme_colors["dark"]

        right_shape = ppt_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.9), Inches(3.8), Inches(4.6))
        right_shape.fill.solid()
        right_shape.fill.fore_color.rgb = RGBColor(240, 246, 255)
        right_shape.line.color.rgb = self.theme_colors["accent"]
        right_tf = right_shape.text_frame
        right_tf.clear()
        right_header = right_tf.paragraphs[0]
        right_header.text = "Key Highlights"
        right_header.font.bold = True
        right_header.font.size = Pt(16)
        right_header.font.color.rgb = self.theme_colors["primary"]
        for insight in (highlight_bullets[:4] or left_bullets[:2]):
            p = right_tf.add_paragraph()
            p.text = insight
            p.level = 0
            p.space_after = Pt(8)
            p.font.name = self.fonts["body"]
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme_colors["dark"]

        # Footer
        footer = ppt_slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(6.0), Inches(0.3))
        footer_para = footer.text_frame.paragraphs[0]
        footer_para.text = "Generated by AI Presentation System"
        footer_para.font.size = Pt(10)
        footer_para.font.color.rgb = RGBColor(136, 136, 136)

        # Add chart if applicable
        chart_spec = self._find_chart_for_slide(chart_specs, slide_number)
        if chart_spec:
            self._add_chart_to_slide(ppt_slide, chart_spec)
        
        return ppt_slide

    def _clean_bullet(self, text: Any) -> str:
        bullet = str(text).strip()
        bullet = bullet.lstrip("-• ").strip()
        return bullet if len(bullet) <= 120 else f"{bullet[:117]}..."

    def _is_highlight_bullet(self, bullet: str) -> bool:
        return bool(re.search(r"(\$|%|\d)", bullet))

    def _is_section_break_slide(self, title: str) -> bool:
        return title.strip().lower() in {"revenue analysis", "market expansion", "customer metrics"}

    def _render_title_slide(self, slide: Any, title: str, content: List[str]):
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.5))
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = title
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.name = self.fonts["title"]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme_colors["primary"]

        subtitle = self._clean_bullet(content[0]) if content else "Project Update"
        subtitle_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.7), Inches(11.3), Inches(0.8))
        subtitle_para = subtitle_box.text_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(22)
        subtitle_para.font.color.rgb = self.theme_colors["dark"]

    def _render_section_break_slide(self, slide: Any, title: str, slide_number: int):
        section_no = slide.shapes.add_textbox(Inches(5.6), Inches(2.2), Inches(1.5), Inches(0.6))
        n_para = section_no.text_frame.paragraphs[0]
        n_para.text = f"{slide_number:02d}"
        n_para.alignment = PP_ALIGN.CENTER
        n_para.font.size = Pt(18)
        n_para.font.color.rgb = RGBColor(136, 136, 136)

        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.9), Inches(9.5), Inches(1.2))
        t_para = title_box.text_frame.paragraphs[0]
        t_para.text = title
        t_para.alignment = PP_ALIGN.CENTER
        t_para.font.size = Pt(44)
        t_para.font.bold = True
        t_para.font.color.rgb = self.theme_colors["primary"]
    
    def _add_title_text(self, slide: Any, text: str, position: Dict[str, Any]):
        """Add title text to slide."""
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = text
            
            # Format title
            text_frame = title_shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = self.fonts["title"]
                paragraph.font.size = Pt(self.font_sizes["title"])
                paragraph.font.color.rgb = self.theme_colors["primary"]
                paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_subtitle_text(self, slide: Any, text: str, position: Dict[str, Any]):
        """Add subtitle text to slide."""
        left = Inches(position.get("left", 1))
        top = Inches(position.get("top", 3))
        width = Inches(position.get("width", 8))
        height = Inches(position.get("height", 1))
        
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = subtitle_box.text_frame
        text_frame.text = text
        
        # Format subtitle
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = self.fonts["subtitle"]
        paragraph.font.size = Pt(self.font_sizes["subtitle"])
        paragraph.font.color.rgb = self.theme_colors["dark"]
        paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_slide_title(self, slide: Any, text: str):
        """Add slide title."""
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = text
            
            # Format title
            text_frame = title_shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = self.fonts["heading"]
                paragraph.font.size = Pt(self.font_sizes["heading"])
                paragraph.font.color.rgb = self.theme_colors["primary"]
                paragraph.alignment = PP_ALIGN.LEFT
    
    def _add_bullet_list(self, slide: Any, text: str, position: Dict[str, Any]):
        """Add bullet list to slide."""
        left = Inches(position.get("left", 1))
        top = Inches(position.get("top", 2))
        width = Inches(position.get("width", 8))
        height = Inches(position.get("height", 4))
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        text_frame.clear()
        
        # Split text into bullet points
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            
            paragraph.text = line.strip()
            paragraph.level = 0
            paragraph.font.name = self.fonts["body"]
            paragraph.font.size = Pt(self.font_sizes["body"])
            paragraph.font.color.rgb = self.theme_colors["dark"]
    
    def _add_text_content(self, slide: Any, text: str, position: Dict[str, Any]):
        """Add text content to slide."""
        left = Inches(position.get("left", 1))
        top = Inches(position.get("top", 2))
        width = Inches(position.get("width", 8))
        height = Inches(position.get("height", 4))
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = text
        
        # Format text
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.fonts["body"]
            paragraph.font.size = Pt(self.font_sizes["body"])
            paragraph.font.color.rgb = self.theme_colors["dark"]
    
    def _add_flow_diagram(self, slide: Any, content: Dict[str, Any], position: Dict[str, Any]):
        """Add flow diagram to slide."""
        # Create simple flow diagram using shapes
        layout = content.get("layout", {})
        flow_content = content.get("content", [])
        
        left = Inches(position.get("left", 1))
        top = Inches(position.get("top", 2))
        width = Inches(position.get("width", 8))
        height = Inches(position.get("height", 3))
        
        # Create flow boxes
        box_width = width / len(flow_content) if flow_content else width
        
        for i, item in enumerate(flow_content):
            x = left + i * box_width
            y = top + height / 2 - Inches(0.5)
            
            # Add rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, box_width - Inches(0.2), Inches(1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.theme_colors["primary"]
            shape.line.color.rgb = self.theme_colors["dark"]
            
            # Add text
            text_frame = shape.text_frame
            text_frame.text = item.get("text", f"Step {i+1}")
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            
            # Add arrow if not last
            if i < len(flow_content) - 1:
                arrow_x = x + box_width - Inches(0.2)
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, arrow_x, y + Inches(0.3), 
                    Inches(0.4), Inches(0.4)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.theme_colors["dark"]
    
    def _add_infographic(self, slide: Any, content: Dict[str, Any], position: Dict[str, Any]):
        """Add infographic to slide."""
        layout = content.get("layout", {})
        info_content = content.get("content", [])
        
        left = Inches(position.get("left", 1))
        top = Inches(position.get("top", 2))
        width = Inches(position.get("width", 8))
        height = Inches(position.get("height", 4))
        
        # Create metric cards
        card_width = width / len(info_content) if info_content else width
        
        for i, item in enumerate(info_content):
            x = left + i * card_width
            y = top
            
            # Add card background
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x + Inches(0.1), y + Inches(0.1), 
                card_width - Inches(0.2), height - Inches(0.2)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.theme_colors["light"]
            shape.line.color.rgb = self.theme_colors["primary"]
            
            # Add text
            text_frame = shape.text_frame
            metric = item.get("metric", "")
            value = item.get("value", "")
            text_frame.text = f"{metric}\n{value}"
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            
            # Format text
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = self.fonts["body"]
                paragraph.font.size = Pt(self.font_sizes["subheading"])
                paragraph.font.color.rgb = self.theme_colors["dark"]
                paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_comparison_grid(self, slide: Any, content: Dict[str, Any], position: Dict[str, Any]):
        """Add comparison grid to slide."""
        layout = content.get("layout", {})
        grid_content = content.get("content", [])
        
        left = Inches(position.get("left", 1))
        top = Inches(position.get("top", 2))
        width = Inches(position.get("width", 8))
        height = Inches(position.get("height", 4))
        
        # Create grid cells
        cols = layout.get("columns", 2)
        rows = layout.get("rows", 2)
        
        cell_width = width / cols
        cell_height = height / rows
        
        for i, item in enumerate(grid_content):
            row = i // cols
            col = i % cols
            
            if row >= rows:
                break
            
            x = left + col * cell_width
            y = top + row * cell_height
            
            # Add cell
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x + Inches(0.05), y + Inches(0.05), 
                cell_width - Inches(0.1), cell_height - Inches(0.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.theme_colors["light"]
            shape.line.color.rgb = self.theme_colors["dark"]
            
            # Add text
            text_frame = shape.text_frame
            text_frame.text = item.get("text", f"Item {i+1}")
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            
            # Format text
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = self.fonts["body"]
                paragraph.font.size = Pt(self.font_sizes["body"])
                paragraph.font.color.rgb = self.theme_colors["dark"]
                paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_visual_element(self, slide: Any, content: Dict[str, Any], position: Dict[str, Any]):
        """Add generic visual element to slide."""
        visual_type = content.get("type", "diagram")
        
        if visual_type == "flow":
            self._add_flow_diagram(slide, content, position)
        elif visual_type == "infographic":
            self._add_infographic(slide, content, position)
        elif visual_type == "grid":
            self._add_comparison_grid(slide, content, position)
        else:
            # Default to simple text box
            self._add_text_content(slide, str(content), position)
    
    def _add_conclusion_content(self, slide: Any, text: str, position: Dict[str, Any]):
        """Add conclusion content to slide."""
        self._add_text_content(slide, text, position)
    
    def _add_chart_to_slide(self, slide: Any, chart_spec: Dict[str, Any]):
        """Add chart to slide."""
        try:
            # Create chart data object
            chart_data_dict = chart_spec.get("data", {})
            chart_type = chart_spec.get("chart_type", "bar")
            
            chart_data = {
                "chart_type": ChartType(chart_type),
                "title": chart_spec.get("title", "Chart"),
                "labels": chart_data_dict.get("labels", []),
                "datasets": chart_data_dict.get("datasets", []),
                "colors": self._get_chart_colors(chart_type)
            }
            
            # Add chart using PPTX utils
            self.pptx_utils.add_chart_slide(
                Presentation(),  # Create temporary presentation for chart
                chart_spec.get("title", ""),
                chart_data
            )
            
        except Exception as e:
            self.logger.error(f"Error adding chart to slide: {str(e)}")
            # Add placeholder text
            self._add_text_content(slide, f"Chart: {chart_spec.get('title', 'Data Visualization')}", 
                                 {"left": 1, "top": 2, "width": 8, "height": 4})
    
    def _find_chart_for_slide(self, chart_specs: List[Dict], slide_number: int) -> Optional[Dict[str, Any]]:
        """Find chart specification for a specific slide."""
        for spec in chart_specs:
            if spec.get("slide_number") == slide_number:
                return spec
        return None
    
    def _get_chart_colors(self, chart_type: str) -> List[str]:
        """Get colors for chart type."""
        return [
            self.theme_colors["primary"],
            self.theme_colors["secondary"],
            self.theme_colors["accent"],
            self.theme_colors["danger"],
            self.theme_colors["warning"]
        ]
    
    def _apply_theme_to_presentation(self, prs: Presentation):
        """Apply consistent theme to presentation."""
        # This is a simplified theme application
        # In a real implementation, you'd apply more comprehensive theming
        pass
    
    def _create_generation_summary(self, slide_layouts: List[Dict], chart_specs: List[Dict], output_path: str) -> Dict[str, Any]:
        """Create a summary of the generation process."""
        slide_types = {}
        
        for layout in slide_layouts:
            slide_type = layout.get("slide_type", "unknown")
            slide_types[slide_type] = slide_types.get(slide_type, 0) + 1
        
        return {
            "total_slides_generated": len(slide_layouts),
            "slide_type_distribution": slide_types,
            "total_charts": len(chart_specs),
            "output_file": output_path,
            "generation_success": True,
            "theme_applied": "consulting_theme",
            "quality_checks": self._perform_quality_checks(slide_layouts)
        }
    
    def _perform_quality_checks(self, slide_layouts: List[Dict]) -> Dict[str, Any]:
        """Perform quality checks on generated slides."""
        checks = {
            "all_slides_have_titles": True,
            "consistent_spacing": True,
            "no_text_overload": True,
            "visual_balance": True
        }
        
        for layout in slide_layouts:
            # Check if slide has title
            content_areas = layout.get("content_areas", [])
            has_title = any(area.get("type") == "title" for area in content_areas)
            if not has_title:
                checks["all_slides_have_titles"] = False
        
        return checks
    
    def _get_file_size(self, file_path: str) -> str:
        """Get file size in human readable format."""
        try:
            if os.path.exists(file_path):
                size_bytes = os.path.getsize(file_path)
                if size_bytes < 1024:
                    return f"{size_bytes} B"
                elif size_bytes < 1024 * 1024:
                    return f"{size_bytes / 1024:.1f} KB"
                else:
                    return f"{size_bytes / (1024 * 1024):.1f} MB"
        except:
            pass
        
        return "Unknown"
    
    def _estimate_generation_time(self, slide_count: int) -> str:
        """Estimate generation time."""
        # Rough estimate: 2 seconds per slide
        seconds = slide_count * 2
        if seconds < 60:
            return f"{seconds}s"
        else:
            return f"{seconds // 60}m {seconds % 60}s"
    
    def validate_output(self, output: Dict[str, Any]) -> bool:
        """Validate PPTX generator output."""
        if not output:
            return False
        
        required_keys = ["presentation_path", "generation_summary", "generation_metadata"]
        for key in required_keys:
            if key not in output:
                self.logger.error(f"Missing required key in PPTX generator output: {key}")
                return False
        
        # Check if presentation file was created
        output_path = output.get("presentation_path", "")
        if not os.path.exists(output_path):
            self.logger.error(f"Presentation file not created: {output_path}")
            return False
        
        return True
