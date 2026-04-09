"""
Slide Master implementation for consistent PowerPoint design.
"""

import os
from typing import Dict, Any, Optional, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class SlideMaster:
    """Manages slide master templates and styling."""
    
    def __init__(self):
        self.theme_colors = {
            "primary": RGBColor(31, 119, 180),      # Blue
            "secondary": RGBColor(255, 127, 14),     # Orange
            "accent": RGBColor(44, 160, 44),         # Green
            "danger": RGBColor(214, 39, 40),         # Red
            "warning": RGBColor(255, 193, 7),        # Yellow
            "info": RGBColor(23, 162, 184),          # Cyan
            "dark": RGBColor(44, 62, 80),            # Dark blue
            "light": RGBColor(236, 240, 241),         # Light gray
            "white": RGBColor(255, 255, 255),         # White
            "gray": RGBColor(128, 128, 128)           # Gray
        }
        
        self.fonts = {
            "title": "Calibri Light",
            "heading": "Calibri",
            "body": "Calibri",
            "subtitle": "Calibri"
        }
        
        self.font_sizes = {
            "title": 44,
            "subtitle": 32,
            "heading": 28,
            "subheading": 20,
            "body": 18,
            "small": 14,
            "caption": 12
        }
        
        # Layout specifications
        self.layouts = self._initialize_layouts()
    
    def apply_theme(self, prs: Presentation, theme_name: str = "consulting") -> Presentation:
        """Apply a theme to the presentation."""
        if theme_name == "consulting":
            return self._apply_consulting_theme(prs)
        elif theme_name == "modern":
            return self._apply_modern_theme(prs)
        elif theme_name == "minimal":
            return self._apply_minimal_theme(prs)
        else:
            return self._apply_default_theme(prs)
    
    def _apply_consulting_theme(self, prs: Presentation) -> Presentation:
        """Apply consulting-style theme."""
        # This is a simplified theme application
        # In a real implementation, you'd modify slide masters, color schemes, etc.
        
        # Set default font for text placeholders
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.font:
                            paragraph.font.name = self.fonts["body"]
                            paragraph.font.size = Pt(self.font_sizes["body"])
        
        return prs
    
    def _apply_modern_theme(self, prs: Presentation) -> Presentation:
        """Apply modern theme."""
        # Modern theme with vibrant colors
        modern_colors = {
            "primary": RGBColor(0, 123, 255),      # Bright blue
            "secondary": RGBColor(255, 85, 0),       # Bright orange
            "accent": RGBColor(0, 200, 83),         # Bright green
        }
        
        # Update theme colors
        self.theme_colors.update(modern_colors)
        
        return prs
    
    def _apply_minimal_theme(self, prs: Presentation) -> Presentation:
        """Apply minimal theme."""
        # Minimal theme with grayscale
        minimal_colors = {
            "primary": RGBColor(33, 33, 33),          # Dark gray
            "secondary": RGBColor(66, 66, 66),        # Medium gray
            "accent": RGBColor(153, 153, 153),        # Light gray
            "light": RGBColor(248, 248, 248),         # Very light gray
        }
        
        # Update theme colors
        self.theme_colors.update(minimal_colors)
        
        return prs
    
    def _apply_default_theme(self, prs: Presentation) -> Presentation:
        """Apply default theme."""
        return prs
    
    def get_layout_specification(self, layout_type: str) -> Dict[str, Any]:
        """Get layout specification for a given type."""
        return self.layouts.get(layout_type, self.layouts["default"])
    
    def _initialize_layouts(self) -> Dict[str, Any]:
        """Initialize layout specifications."""
        return {
            "title": {
                "name": "Title Slide",
                "description": "Main title slide with subtitle",
                "placeholders": {
                    "title": {
                        "position": {"left": 1.0, "top": 2.0},
                        "size": {"width": 8.0, "height": 1.5},
                        "font": self.fonts["title"],
                        "font_size": self.font_sizes["title"],
                        "color": self.theme_colors["primary"],
                        "alignment": PP_ALIGN.CENTER
                    },
                    "subtitle": {
                        "position": {"left": 1.5, "top": 4.0},
                        "size": {"width": 7.0, "height": 1.0},
                        "font": self.fonts["subtitle"],
                        "font_size": self.font_sizes["subtitle"],
                        "color": self.theme_colors["dark"],
                        "alignment": PP_ALIGN.CENTER
                    }
                },
                "background": self.theme_colors["white"]
            },
            "content": {
                "name": "Content Slide",
                "description": "Standard content slide with title and bullet points",
                "placeholders": {
                    "title": {
                        "position": {"left": 1.0, "top": 0.5},
                        "size": {"width": 8.0, "height": 1.0},
                        "font": self.fonts["heading"],
                        "font_size": self.font_sizes["heading"],
                        "color": self.theme_colors["primary"],
                        "alignment": PP_ALIGN.LEFT
                    },
                    "content": {
                        "position": {"left": 1.0, "top": 2.0},
                        "size": {"width": 8.0, "height": 4.5},
                        "font": self.fonts["body"],
                        "font_size": self.font_sizes["body"],
                        "color": self.theme_colors["dark"],
                        "alignment": PP_ALIGN.LEFT
                    }
                },
                "background": self.theme_colors["white"]
            },
            "two_column": {
                "name": "Two Column Slide",
                "description": "Slide with two columns of content",
                "placeholders": {
                    "title": {
                        "position": {"left": 1.0, "top": 0.5},
                        "size": {"width": 8.0, "height": 1.0},
                        "font": self.fonts["heading"],
                        "font_size": self.font_sizes["heading"],
                        "color": self.theme_colors["primary"],
                        "alignment": PP_ALIGN.LEFT
                    },
                    "left_content": {
                        "position": {"left": 1.0, "top": 2.0},
                        "size": {"width": 3.5, "height": 4.5},
                        "font": self.fonts["body"],
                        "font_size": self.font_sizes["body"],
                        "color": self.theme_colors["dark"],
                        "alignment": PP_ALIGN.LEFT
                    },
                    "right_content": {
                        "position": {"left": 5.0, "top": 2.0},
                        "size": {"width": 3.5, "height": 4.5},
                        "font": self.fonts["body"],
                        "font_size": self.font_sizes["body"],
                        "color": self.theme_colors["dark"],
                        "alignment": PP_ALIGN.LEFT
                    }
                },
                "background": self.theme_colors["white"]
            },
            "section": {
                "name": "Section Divider",
                "description": "Section divider slide",
                "placeholders": {
                    "title": {
                        "position": {"left": 1.0, "top": 3.0},
                        "size": {"width": 8.0, "height": 1.5},
                        "font": self.fonts["heading"],
                        "font_size": self.font_sizes["heading"],
                        "color": self.theme_colors["primary"],
                        "alignment": PP_ALIGN.CENTER
                    }
                },
                "background": self.theme_colors["light"]
            },
            "chart": {
                "name": "Chart Slide",
                "description": "Slide with chart and optional content",
                "placeholders": {
                    "title": {
                        "position": {"left": 1.0, "top": 0.5},
                        "size": {"width": 8.0, "height": 1.0},
                        "font": self.fonts["heading"],
                        "font_size": self.font_sizes["heading"],
                        "color": self.theme_colors["primary"],
                        "alignment": PP_ALIGN.LEFT
                    },
                    "chart": {
                        "position": {"left": 1.0, "top": 2.0},
                        "size": {"width": 6.0, "height": 4.0},
                        "background": self.theme_colors["white"]
                    },
                    "content": {
                        "position": {"left": 7.5, "top": 2.0},
                        "size": {"width": 1.5, "height": 4.0},
                        "font": self.fonts["body"],
                        "font_size": self.font_sizes["small"],
                        "color": self.theme_colors["dark"],
                        "alignment": PP_ALIGN.LEFT
                    }
                },
                "background": self.theme_colors["white"]
            },
            "default": {
                "name": "Default Layout",
                "description": "Default slide layout",
                "placeholders": {
                    "title": {
                        "position": {"left": 1.0, "top": 0.5},
                        "size": {"width": 8.0, "height": 1.0},
                        "font": self.fonts["heading"],
                        "font_size": self.font_sizes["heading"],
                        "color": self.theme_colors["primary"],
                        "alignment": PP_ALIGN.LEFT
                    },
                    "content": {
                        "position": {"left": 1.0, "top": 2.0},
                        "size": {"width": 8.0, "height": 4.5},
                        "font": self.fonts["body"],
                        "font_size": self.font_sizes["body"],
                        "color": self.theme_colors["dark"],
                        "alignment": PP_ALIGN.LEFT
                    }
                },
                "background": self.theme_colors["white"]
            }
        }
    
    def create_custom_template(self, name: str, specifications: Dict[str, Any]) -> Dict[str, Any]:
        """Create a custom template based on specifications."""
        custom_template = {
            "name": name,
            "description": specifications.get("description", f"Custom template: {name}"),
            "placeholders": specifications.get("placeholders", {}),
            "background": specifications.get("background", self.theme_colors["white"])
        }
        
        self.layouts[name] = custom_template
        return custom_template
    
    def export_template(self, template_name: str, output_path: str) -> bool:
        """Export a template to a file."""
        try:
            template = self.layouts.get(template_name)
            if not template:
                return False
            
            import json
            with open(output_path, 'w') as f:
                json.dump(template, f, indent=2, default=str)
            
            return True
        except Exception:
            return False
    
    def import_template(self, template_path: str) -> Optional[Dict[str, Any]]:
        """Import a template from a file."""
        try:
            import json
            with open(template_path, 'r') as f:
                template = json.load(f)
            
            # Validate template structure
            if all(key in template for key in ["name", "placeholders", "background"]):
                template_name = template["name"]
                self.layouts[template_name] = template
                return template
            
            return None
        except Exception:
            return None
    
    def list_available_templates(self) -> List[str]:
        """List all available templates."""
        return list(self.layouts.keys())
    
    def get_theme_colors(self) -> Dict[str, RGBColor]:
        """Get current theme colors."""
        return self.theme_colors.copy()
    
    def get_fonts(self) -> Dict[str, str]:
        """Get current font settings."""
        return self.fonts.copy()
    
    def get_font_sizes(self) -> Dict[str, int]:
        """Get current font size settings."""
        return self.font_sizes.copy()
